"""Document parsing service for different file types."""

import io
import re
import logging
from pathlib import Path
from typing import Optional
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    """Parsed document content."""
    filename: str
    file_type: str
    text: str
    pages: list[str]
    tables: list[dict]
    metadata: dict


class DocumentParser:
    """Parses various document formats to extract text and structure."""

    SUPPORTED_TYPES = {
        ".pdf": "pdf",
        ".xlsx": "excel",
        ".xls": "excel",
        ".docx": "word",
        ".doc": "word",
        ".pptx": "powerpoint",
        ".ppt": "powerpoint",
        ".csv": "csv",
        ".txt": "text",
    }

    def parse(self, file_path: str | Path, content: Optional[bytes] = None) -> ParsedDocument:
        """Parse a document and extract text content."""
        path = Path(file_path)
        ext = path.suffix.lower()
        file_type = self.SUPPORTED_TYPES.get(ext, "unknown")

        if content is None:
            with open(path, "rb") as f:
                content = f.read()

        if file_type == "pdf":
            return self._parse_pdf(path.name, content)
        elif file_type == "excel":
            return self._parse_excel(path.name, content)
        elif file_type == "word":
            return self._parse_word(path.name, content)
        elif file_type == "powerpoint":
            return self._parse_powerpoint(path.name, content)
        elif file_type == "csv":
            return self._parse_csv(path.name, content)
        elif file_type == "text":
            return self._parse_text(path.name, content)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _parse_pdf(self, filename: str, content: bytes) -> ParsedDocument:
        """Parse PDF document using PyMuPDF."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF not installed, using mock parser")
            return self._mock_parse(filename, "pdf")

        try:
            doc = fitz.open(stream=content, filetype="pdf")
            pages = []
            tables = []
            full_text = []
            
            logger.info(f"Parsing PDF: {filename} ({len(doc)} pages)")

            for page_num, page in enumerate(doc):
                # Extract text with layout preservation
                text = page.get_text("text")
                
                # Also try blocks for better structure
                blocks = page.get_text("blocks")
                block_texts = []
                for block in blocks:
                    if block[6] == 0:  # Text block (not image)
                        block_texts.append(block[4])
                
                # Use block text if available, otherwise use page text
                page_text = "\n".join(block_texts) if block_texts else text
                
                pages.append(page_text)
                full_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
                
                logger.debug(f"Page {page_num + 1}: Extracted {len(page_text)} characters")

                # Extract tables using PyMuPDF's table finder
                try:
                    tabs = page.find_tables()
                    logger.info(f"Page {page_num + 1}: Found {len(tabs)} tables")
                    
                    for tab_idx, tab in enumerate(tabs):
                        try:
                            table_data = tab.extract()
                            if table_data:
                                tables.append({
                                    "page": page_num + 1,
                                    "table_index": tab_idx,
                                    "data": table_data,
                                    "bbox": list(tab.bbox) if hasattr(tab, 'bbox') else None
                                })
                                logger.debug(f"Page {page_num + 1}, Table {tab_idx}: {len(table_data)} rows")
                        except Exception as e:
                            logger.warning(f"Error extracting table {tab_idx} from page {page_num + 1}: {e}")
                            continue
                except Exception as e:
                    logger.warning(f"Error finding tables on page {page_num + 1}: {e}")

            doc.close()
            
            full_text_combined = "\n".join(full_text)
            logger.info(f"PDF parsing complete: {len(full_text_combined)} total characters, {len(tables)} tables")
            
            # Log first 500 characters for debugging
            preview = full_text_combined[:500].replace('\n', ' ')
            logger.debug(f"Extracted text preview: {preview}...")

            return ParsedDocument(
                filename=filename,
                file_type="pdf",
                text=full_text_combined,
                pages=pages,
                tables=tables,
                metadata={
                    "page_count": len(pages),
                    "table_count": len(tables),
                    "total_characters": len(full_text_combined),
                    "parser": "PyMuPDF"
                }
            )
        except Exception as e:
            logger.error(f"Error parsing PDF {filename}: {e}", exc_info=True)
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    def _parse_excel(self, filename: str, content: bytes) -> ParsedDocument:
        """Parse Excel document."""
        try:
            import pandas as pd
        except ImportError:
            return self._mock_parse(filename, "excel")

        xlsx = pd.ExcelFile(io.BytesIO(content))
        pages = []
        tables = []
        full_text = []

        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            text = f"--- Sheet: {sheet_name} ---\n"
            text += df.to_string()
            pages.append(text)
            full_text.append(text)

            tables.append({
                "sheet": sheet_name,
                "headers": df.columns.tolist(),
                "data": df.values.tolist()
            })

        return ParsedDocument(
            filename=filename,
            file_type="excel",
            text="\n".join(full_text),
            pages=pages,
            tables=tables,
            metadata={"sheet_count": len(xlsx.sheet_names)}
        )

    def _parse_word(self, filename: str, content: bytes) -> ParsedDocument:
        """Parse Word document."""
        try:
            from docx import Document
        except ImportError:
            return self._mock_parse(filename, "word")

        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []

        for table in doc.tables:
            table_data = []
            for row in table.rows:
                table_data.append([cell.text for cell in row.cells])
            tables.append({"data": table_data})

        return ParsedDocument(
            filename=filename,
            file_type="word",
            text="\n".join(paragraphs),
            pages=["\n".join(paragraphs)],
            tables=tables,
            metadata={"paragraph_count": len(paragraphs)}
        )

    def _parse_powerpoint(self, filename: str, content: bytes) -> ParsedDocument:
        """Parse PowerPoint document."""
        try:
            from pptx import Presentation
        except ImportError:
            return self._mock_parse(filename, "powerpoint")

        prs = Presentation(io.BytesIO(content))
        pages = []
        tables = []
        full_text = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {slide_num + 1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    tables.append({"slide": slide_num + 1, "data": table_data})

            text = "\n".join(slide_text)
            pages.append(text)
            full_text.append(text)

        return ParsedDocument(
            filename=filename,
            file_type="powerpoint",
            text="\n".join(full_text),
            pages=pages,
            tables=tables,
            metadata={"slide_count": len(prs.slides)}
        )

    def _parse_csv(self, filename: str, content: bytes) -> ParsedDocument:
        """Parse CSV document."""
        try:
            import pandas as pd
        except ImportError:
            return self._mock_parse(filename, "csv")

        df = pd.read_csv(io.BytesIO(content))
        text = df.to_string()

        return ParsedDocument(
            filename=filename,
            file_type="csv",
            text=text,
            pages=[text],
            tables=[{"headers": df.columns.tolist(), "data": df.values.tolist()}],
            metadata={"row_count": len(df)}
        )

    def _parse_text(self, filename: str, content: bytes) -> ParsedDocument:
        """Parse plain text document."""
        text = content.decode("utf-8", errors="ignore")
        return ParsedDocument(
            filename=filename,
            file_type="text",
            text=text,
            pages=[text],
            tables=[],
            metadata={}
        )

    def _mock_parse(self, filename: str, file_type: str) -> ParsedDocument:
        """Return mock data when parser libraries are not installed."""
        return ParsedDocument(
            filename=filename,
            file_type=file_type,
            text=f"[Mock parsed content for {filename}]",
            pages=["[Mock page content]"],
            tables=[],
            metadata={"mock": True}
        )
